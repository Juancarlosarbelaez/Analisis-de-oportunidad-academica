"""
Análisis de Tendencias de Mercado con Agentes de IA (versión con strands)

Este script implementa una arquitectura basada en agentes de inteligencia
artificial para analizar tendencias de mercado en el ámbito académico.  A
partir de la información contenida en la tabla ``programas_equivalentes.xlsx``,
un agente de análisis carga y preprocesa los datos, un agente de tendencias
identifica patrones y un agente de visualización genera gráficas y
tableros.  Además, un agente de recomendaciones sintetiza insights y
conclusiones estratégicas.

Para acercarnos al patrón "planner–executor" y al concepto de ``strands``
mostrado en el archivo de ejemplo, las fases de identificación de tendencias
y generación de visualizaciones se ejecutan de forma concurrente mediante el
módulo ``asyncio``.  Cada subanálisis o gráfica se lanza en un hilo
independiente a través de ``asyncio.to_thread``, de manera que los
"strands" o hebras de ejecución colaboran como agentes autónomos que
contribuyen a un objetivo común.
"""


import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
import numpy as np
from datetime import datetime
import warnings
warnings.filterwarnings('ignore')

# -----------------------------------------------------------------------------
# Integración con la arquitectura de strands
# -----------------------------------------------------------------------------
# Importamos la clase Agent y el decorador tool desde nuestro stub local
from strands import Agent, tool
from model_config import get_configured_model


# Para permitir la ejecución concurrente de distintas tareas de análisis y
# visualización empleamos el módulo asyncio.  Aunque el ejemplo proporcionado
# utiliza la biblioteca `strands` (no disponible en este entorno), la idea
# de «strands» se refiere a ejecutar subtareas en paralelo por medio de
# agentes cooperantes.  Con `asyncio` podemos simular ese patrón de agentes
# cooperando de forma asíncrona: definimos corrutinas que delegan el
# procesamiento a hilos mediante `asyncio.to_thread`. Esto resulta útil para
# lanzar análisis y creación de gráficas en paralelo sin bloquear la
# ejecución principal.
import asyncio
from agents import (
    AgenteAnalisisDatos,
    AgenteTendenciasMercado,
    AgenteVisualizacion,
    AgenteRecomendaciones,
    AgentePresentacion,
)
 


class AgenteRecomendaciones:
    """
    Agente de IA especializado en generar insights y recomendaciones.
    Analiza los resultados y proporciona conclusiones estratégicas.
    """
    
    def __init__(self, datos, tendencias):
        self.datos = datos
        self.tendencias = tendencias
        self.recomendaciones = []
        
    def generar_insights(self):
        """Genera insights basados en el análisis de datos"""
        print("\n🤖 Agente de Recomendaciones: Generando insights...")
        
        insights = []
        
        # Insight 1: Institución líder
        if 'institucional' in self.tendencias and len(self.tendencias['institucional']) > 0:
            institucion_lider = self.tendencias['institucional'].index[0]
            matricula_lider = self.tendencias['institucional'].values[0]
            insights.append(f"📌 La institución líder en matrícula es '{institucion_lider}' "
                          f"con {int(matricula_lider):,} estudiantes matriculados.")
        
        # Insight 2: Programa más popular
        if 'programas' in self.tendencias and len(self.tendencias['programas']) > 0:
            programa_popular = self.tendencias['programas'].index[0]
            matricula_programa = self.tendencias['programas'].values[0]
            insights.append(f"📌 El programa académico más popular es '{programa_popular}' "
                          f"con {int(matricula_programa):,} estudiantes.")
        
        # Insight 3: Distribución sectorial
        if 'sectores' in self.tendencias and len(self.tendencias['sectores']) > 0:
            total_matricula = self.tendencias['sectores'].sum()
            for sector, matricula in self.tendencias['sectores'].items():
                porcentaje = (matricula / total_matricula) * 100
                insights.append(f"📌 El sector {sector} representa el {porcentaje:.1f}% "
                              f"de la matrícula total.")
        
        # Insight 4: Área de conocimiento dominante
        if 'areas' in self.tendencias and len(self.tendencias['areas']) > 0:
            area_dominante = self.tendencias['areas'].index[0]
            matricula_area = self.tendencias['areas'].values[0]
            insights.append(f"📌 El área de conocimiento dominante es '{area_dominante}' "
                          f"con {int(matricula_area):,} estudiantes.")
        
        return insights
    
    def generar_recomendaciones_estrategicas(self):
        """Genera recomendaciones estratégicas"""
        print("\n🤖 Agente de Recomendaciones: Generando recomendaciones estratégicas...")
        
        recomendaciones = []
        
        # Recomendación 1: Oportunidades de mercado
        if 'programas' in self.tendencias and len(self.tendencias['programas']) > 0:
            recomendaciones.append("💡 Oportunidad de Mercado: Los programas con mayor demanda "
                                 "representan áreas de alta oportunidad para inversión y expansión.")
        
        # Recomendación 2: Análisis competitivo
        if 'institucional' in self.tendencias and len(self.tendencias['institucional']) > 0:
            recomendaciones.append("💡 Análisis Competitivo: Las instituciones líderes marcan "
                                 "las tendencias del mercado. Estudiar sus estrategias puede "
                                 "proporcionar ventajas competitivas.")
        
        # Recomendación 3: Diversificación
        if 'areas' in self.tendencias and len(self.tendencias['areas']) > 1:
            recomendaciones.append("💡 Diversificación: Considerar la diversificación hacia "
                                 "múltiples áreas de conocimiento para mitigar riesgos y "
                                 "capturar diferentes segmentos de mercado.")
        
        # Recomendación 4: Sector público vs privado
        if 'sectores' in self.tendencias and len(self.tendencias['sectores']) > 0:
            recomendaciones.append("💡 Estrategia Sectorial: Analizar las diferencias entre "
                                 "sectores público y privado para identificar nichos "
                                 "específicos de oportunidad.")
        
        return recomendaciones
    
    def generar_reporte_completo(self, insights, recomendaciones):
        """Genera un reporte completo en formato texto"""
        print("\n🤖 Agente de Recomendaciones: Generando reporte completo...")
        
        reporte = []
        reporte.append("=" * 80)
        reporte.append("REPORTE DE ANÁLISIS DE TENDENCIAS DE MERCADO ACADÉMICO")
        reporte.append("=" * 80)
        reporte.append(f"Fecha de generación: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
        reporte.append("")
        
        reporte.append("\n" + "=" * 80)
        reporte.append("INSIGHTS PRINCIPALES")
        reporte.append("=" * 80)
        for insight in insights:
            reporte.append(insight)
        
        reporte.append("\n" + "=" * 80)
        reporte.append("RECOMENDACIONES ESTRATÉGICAS")
        reporte.append("=" * 80)
        for rec in recomendaciones:
            reporte.append(rec)
        
        reporte.append("\n" + "=" * 80)
        reporte.append("CONCLUSIÓN")
        reporte.append("=" * 80)
        reporte.append("Este análisis proporciona una visión integral de las tendencias de mercado")
        reporte.append("en el sector académico. Las visualizaciones generadas permiten identificar")
        reporte.append("patrones clave, oportunidades de crecimiento y áreas de mejora estratégica.")
        reporte.append("=" * 80)
        
        # Guardar reporte
        with open('reporte_analisis_tendencias.txt', 'w', encoding='utf-8') as f:
            f.write('\n'.join(reporte))
        
        print("✓ Reporte guardado: reporte_analisis_tendencias.txt")
        
        # También imprimir en consola
        print('\n'.join(reporte))

from pptx import Presentation
from pptx.util import Inches

class AgentePresentacion:
    def __init__(self, resumen, tendencias, insights, recomendaciones, imagenes):
        self.resumen = resumen
        self.tendencias = tendencias
        self.insights = insights
        self.recomendaciones = recomendaciones
        self.imagenes = imagenes  # lista de rutas de imágenes (png)

    def generar_presentacion(self, nombre_archivo="presentacion_analisis_tendencias.pptx"):
        prs = Presentation()

        # 1. Diapositiva de título
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Análisis de Tendencias de Mercado Académico"
        slide.placeholders[1].text = "Reporte ejecutivo generado por agentes de IA"

        # 2. Diapositiva de resumen
        summary_slide = prs.slides.add_slide(prs.slide_layouts[1])
        summary_slide.shapes.title.text = "Resumen de datos"
        cuerpo = summary_slide.placeholders[1].text_frame
        cuerpo.clear()
        cuerpo.add_paragraph().text = f"Total de registros: {self.resumen['total_registros']:,}"
        cuerpo.add_paragraph().text = f"Instituciones únicas: {self.resumen['instituciones']:,}"
        cuerpo.add_paragraph().text = f"Programas académicos únicos: {self.resumen['programas']:,}"
        cuerpo.add_paragraph().text = f"Períodos analizados: {self.resumen['periodos']:,}"

        # 3. Diapositivas con gráficos
        for ruta in self.imagenes:
            titulo = ruta.replace(".png", "").replace("_", " ").title()
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = titulo
            slide.shapes.add_picture(ruta, Inches(0.5), Inches(1.5), width=Inches(9))

        # 4. Diapositiva de insights y recomendaciones
        info_slide = prs.slides.add_slide(prs.slide_layouts[1])
        info_slide.shapes.title.text = "Insights y recomendaciones"
        txt_frame = info_slide.placeholders[1].text_frame
        txt_frame.clear()
        txt_frame.add_paragraph().text = "📌 Principales insights:"
        for insight in self.insights:
            p = txt_frame.add_paragraph()
            p.text = insight
            p.level = 1
        txt_frame.add_paragraph().text = "💡 Recomendaciones estratégicas:"
        for rec in self.recomendaciones:
            p = txt_frame.add_paragraph()
            p.text = rec
            p.level = 1

        # 5. Guardar presentación
        prs.save(nombre_archivo)
        print(f"✓ Presentación guardada: {nombre_archivo}")


def main():
    """
    Punto de entrada de la aplicación.

    En lugar de ejecutar las fases de manera secuencial y explícita, aquí se
    configura un agente tipo Planner de acuerdo con la arquitectura
    ``strands``.  El Planner delega cada fase del proceso a herramientas
    registradas mediante el decorador ``@tool``.  Estas herramientas
    actualizan un contexto compartido y generan los artefactos (gráficas,
    reportes, presentación) conforme se ejecutan.
    """
    # Contenedor de contexto global donde las herramientas almacenarán sus
    # resultados.  Esto incluye los datos, tendencias, resumen, listas de
    # imágenes, insights y recomendaciones.
    global _context
    _context = {}

    # Ruta al archivo de datos que las herramientas necesitan conocer.
    archivo_excel = 'programas_equivalentes.xlsx'

    # ------------------------------------------------------------------
    # Definición de herramientas específicas para cada fase del pipeline.
    # Utilizamos el decorador @tool para que sean compatibles con la
    # infraestructura de strands.  Cada herramienta toma un prompt como
    # argumento pero no lo utiliza directamente; en su lugar, realiza su
    # tarea y actualiza el contexto global.
    # ------------------------------------------------------------------

    @tool
    def herramienta_carga_y_preprocesamiento(prompt: str) -> str:
        """Carga y preprocesa los datos del archivo Excel.

        Actualiza ``_context['datos']`` con el DataFrame resultante y
        ``_context['resumen']`` con el resumen estadístico.  Devuelve un
        mensaje indicando el éxito de la operación.
        """
        print("\n📋 FASE 1: ANÁLISIS Y PREPROCESAMIENTO DE DATOS")
        print("-" * 80)
        agente_datos = AgenteAnalisisDatos(archivo_excel)
        datos = agente_datos.cargar_datos()
        datos = agente_datos.preprocesar_datos()
        resumen = agente_datos.obtener_resumen_estadistico()
        # Guardar en contexto
        _context['datos'] = datos
        _context['resumen'] = resumen
        return "✓ Datos cargados y preprocesados."

    @tool
    def herramienta_identificacion_tendencias(prompt: str) -> str:
        """Identifica las tendencias de mercado a partir de los datos.

        Instancia un ``AgenteTendenciasMercado`` y ejecuta todos los
        análisis de forma concurrente mediante ``asyncio.run``.  Al final
        guarda el diccionario de tendencias en ``_context['tendencias']``.
        """
        print("\n📊 FASE 2: IDENTIFICACIÓN DE TENDENCIAS DE MERCADO")
        print("-" * 80)
        datos = _context.get('datos')
        if datos is None:
            return "⚠️ No se encontraron datos para analizar."
        agente_tendencias = AgenteTendenciasMercado(datos)
        # Ejecutar todas las tendencias en paralelo.  Cualquier excepción
        # se propagará y será visible en la consola.
        asyncio.run(agente_tendencias.analizar_todo_concurrente())
        # Guardar resultados
        _context['tendencias'] = agente_tendencias.tendencias
        return "✓ Tendencias identificadas."

    @tool
    def herramienta_generacion_visualizaciones(prompt: str) -> str:
        """Genera las visualizaciones y el dashboard completos.

        Utiliza ``AgenteVisualizacion`` para crear todas las gráficas de
        forma concurrente.  Guarda las rutas de las imágenes en
        ``_context['imagenes']``.
        """
        print("\n📈 FASE 3: GENERACIÓN DE VISUALIZACIONES")
        print("-" * 80)
        datos = _context.get('datos')
        tendencias = _context.get('tendencias')
        if datos is None or tendencias is None:
            return "⚠️ No se pueden generar visualizaciones por falta de datos o tendencias."
        agente_viz = AgenteVisualizacion(datos, tendencias)
        asyncio.run(agente_viz.crear_todas_graficas(top_n=10))
        _context['imagenes'] = agente_viz.figuras
        return f"✓ Visualizaciones generadas ({len(agente_viz.figuras)} archivos)."

    @tool
    def herramienta_insights_y_recomendaciones(prompt: str) -> str:
        """Genera insights, recomendaciones y el reporte de texto.

        Instancia ``AgenteRecomendaciones`` para producir los insights y
        recomendaciones.  También se genera un reporte de texto y se
        actualizan las claves ``_context['insights']`` y
        ``_context['recomendaciones']``.
        """
        print("\n💡 FASE 4: GENERACIÓN DE INSIGHTS Y RECOMENDACIONES")
        print("-" * 80)
        datos = _context.get('datos')
        tendencias = _context.get('tendencias')
        if datos is None or tendencias is None:
            return "⚠️ No se pueden generar insights debido a datos incompletos."
        agente_recom = AgenteRecomendaciones(datos, tendencias)
        insights = agente_recom.generar_insights()
        recomendaciones = agente_recom.generar_recomendaciones_estrategicas()
        agente_recom.generar_reporte_completo(insights, recomendaciones)
        # Guardar en contexto
        _context['insights'] = insights
        _context['recomendaciones'] = recomendaciones
        return f"✓ Insights y recomendaciones generados ({len(insights)} insights, {len(recomendaciones)} recomendaciones)."

    @tool
    def herramienta_generar_presentacion(prompt: str) -> str:
        """Genera la presentación en PowerPoint basada en los resultados.

        Usa ``AgentePresentacion`` para armar una presentación con el
        resumen de datos, tendencias, insights, recomendaciones e
        imágenes.  El nombre del archivo resultante es
        ``presentacion_analisis_tendencias.pptx``.
        """
        print("\n🎞️ FASE 5: GENERACIÓN DE PRESENTACIÓN")
        print("-" * 80)
        resumen = _context.get('resumen')
        tendencias = _context.get('tendencias')
        insights = _context.get('insights')
        recomendaciones = _context.get('recomendaciones')
        imagenes = _context.get('imagenes')
        if not all([resumen, tendencias, insights, recomendaciones, imagenes]):
            return "⚠️ La presentación no se puede generar por información incompleta."
        agente_presentacion = AgentePresentacion(
            resumen=resumen,
            tendencias=tendencias,
            insights=insights,
            recomendaciones=recomendaciones,
            imagenes=imagenes
        )
        agente_presentacion.generar_presentacion("presentacion_analisis_tendencias.pptx")
        return "✓ Presentación generada."

    # ------------------------------------------------------------------
    # Configuración del planificador (Planner)
    # ------------------------------------------------------------------
    # El Planner se configura con las herramientas definidas arriba.  El
    # sistema_prompt aquí es informativo, ya que nuestra implementación
    # simplificada de Agent no lo utiliza para planificar.  Aun así, se
    # incluye para mantener paralelismo con los ejemplos originales.
    planner_instructions = (
        "Eres un PLANNER responsable de coordinar un análisis de tendencias "
        "de mercado académico.  Deberás invocar las herramientas en el "
        "orden correcto: (1) carga y preprocesamiento de datos, (2) "
        "identificación de tendencias, (3) generación de visualizaciones, "
        "(4) insights y recomendaciones, y (5) creación de la presentación."
    )

    planner = Agent(
        model=get_configured_model(),
        name="Planner",
        system_prompt=planner_instructions,
        tools=[
            herramienta_carga_y_preprocesamiento,
            herramienta_identificacion_tendencias,
            herramienta_generacion_visualizaciones,
            herramienta_insights_y_recomendaciones,
            herramienta_generar_presentacion,
        ]
    )

    # ------------------------------------------------------------------
    # Ejecución del planificador
    # ------------------------------------------------------------------
    print("\n" + "=" * 80)
    print("SISTEMA DE ANÁLISIS DE TENDENCIAS DE MERCADO CON AGENTES DE IA")
    print("=" * 80)
    print("Iniciando ejecución a través del Planner de strands...")
    print("=" * 80 + "\n")

    # Al invocar el Planner con un prompt, se ejecutarán en secuencia
    # todas las herramientas registradas.  La salida de cada herramienta se
    # imprimirá por pantalla y se concatenará en la respuesta final.
    resultado = planner("Iniciar análisis de tendencias de mercado académico")
    # El objeto devuelto por `planner` puede no ser un string (p. ej. AgentResult).
    # Convertimos a cadena para evitar TypeError al concatenar.
    print("\n" + str(resultado))

    # Resumen final de archivos generados
    imagenes = _context.get('imagenes', [])
    print("\n" + "=" * 80)
    print("✅ ANÁLISIS COMPLETADO EXITOSAMENTE")
    print("=" * 80)
    print(f"\nArchivos generados:")
    for i, figura in enumerate(imagenes, 1):
        print(f"  {i}. {figura}")
    print(f"  {len(imagenes) + 1}. reporte_analisis_tendencias.txt")
    print(f"  {len(imagenes) + 2}. presentacion_analisis_tendencias.pptx")
    print("\n" + "=" * 80 + "\n")


if __name__ == "__main__":
    main()