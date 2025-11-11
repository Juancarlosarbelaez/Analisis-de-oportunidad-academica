"""
Análisis de Tendencias de Mercado con Agentes de IA
Este script utiliza agentes de inteligencia artificial para analizar las tendencias
de mercado académico y presentar los resultados mediante gráficas, a partir de la
información contenida en la tabla "programas_equivalentes.xlsx"
"""

import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
import numpy as np
from datetime import datetime
import warnings
warnings.filterwarnings('ignore')

# Configuración de estilo para las gráficas
plt.style.use('seaborn-v0_8-darkgrid')
sns.set_palette("husl")

class AgenteAnalisisDatos:
    """
    Agente de IA especializado en análisis de datos académicos.
    Procesa y limpia los datos para prepararlos para el análisis.
    """
    
    def __init__(self, archivo_excel):
        self.archivo = archivo_excel
        self.datos = None
        
    def cargar_datos(self):
        """Carga los datos del archivo Excel"""
        print("🤖 Agente de Análisis: Cargando datos...")
        self.datos = pd.read_excel(self.archivo)
        print(f"✓ Datos cargados: {self.datos.shape[0]} registros, {self.datos.shape[1]} columnas")
        return self.datos
    
    def preprocesar_datos(self):
        """Preprocesa y limpia los datos"""
        print("\n🤖 Agente de Análisis: Preprocesando datos...")
        
        # Limpiar valores nulos en columnas importantes
        self.datos['MATRICULA'].fillna(0, inplace=True)
        
        # Crear columnas derivadas útiles
        self.datos['PERIODO_COMPLETO'] = self.datos['AÑO_x'].astype(str) + '-' + self.datos['SEMESTRE_x'].astype(str)
        
        print("✓ Datos preprocesados")
        return self.datos
    
    def obtener_resumen_estadistico(self):
        """Genera un resumen estadístico de los datos"""
        print("\n📊 Resumen Estadístico:")
        print(f"   • Total de registros: {len(self.datos)}")
        print(f"   • Instituciones únicas: {self.datos['INSTITUCION'].nunique()}")
        print(f"   • Programas académicos únicos: {self.datos['PROGRAMA_ACADEMICO'].nunique()}")
        print(f"   • Períodos analizados: {self.datos['PERIODO'].nunique()}")
        return {
            'total_registros': len(self.datos),
            'instituciones': self.datos['INSTITUCION'].nunique(),
            'programas': self.datos['PROGRAMA_ACADEMICO'].nunique(),
            'periodos': self.datos['PERIODO'].nunique()
        }


class AgenteTendenciasMercado:
    """
    Agente de IA especializado en identificación de tendencias de mercado.
    Analiza patrones temporales y de crecimiento en la matrícula.
    """
    
    def __init__(self, datos):
        self.datos = datos
        self.tendencias = {}
        
    def analizar_tendencia_matricula_temporal(self):
        """Analiza la tendencia de matrícula a lo largo del tiempo"""
        print("\n🤖 Agente de Tendencias: Analizando evolución temporal de matrículas...")
        
        # Agrupar por período y calcular matrícula total
        tendencia_temporal = self.datos.groupby('PERIODO')['MATRICULA'].sum().sort_index()
        
        self.tendencias['temporal'] = tendencia_temporal
        print(f"✓ Identificados {len(tendencia_temporal)} períodos con datos de matrícula")
        return tendencia_temporal
    
    def analizar_tendencia_por_institucion(self):
        """Analiza tendencias por institución"""
        print("\n🤖 Agente de Tendencias: Analizando tendencias por institución...")
        
        # Top instituciones por matrícula total
        tendencia_institucional = self.datos.groupby('INSTITUCION')['MATRICULA'].sum().sort_values(ascending=False)
        
        self.tendencias['institucional'] = tendencia_institucional
        print(f"✓ Analizadas {len(tendencia_institucional)} instituciones")
        return tendencia_institucional
    
    def analizar_tendencia_por_programa(self):
        """Analiza tendencias por programa académico"""
        print("\n🤖 Agente de Tendencias: Analizando tendencias por programa académico...")
        
        tendencia_programas = self.datos.groupby('PROGRAMA_ACADEMICO')['MATRICULA'].sum().sort_values(ascending=False)
        
        self.tendencias['programas'] = tendencia_programas
        print(f"✓ Analizados {len(tendencia_programas)} programas académicos")
        return tendencia_programas
    
    def analizar_distribucion_sectores(self):
        """Analiza la distribución entre sectores público y privado"""
        print("\n🤖 Agente de Tendencias: Analizando distribución por sector...")
        
        tendencia_sectores = self.datos.groupby('SECTOR_IES')['MATRICULA'].sum()
        
        self.tendencias['sectores'] = tendencia_sectores
        print(f"✓ Analizados sectores: {', '.join(tendencia_sectores.index.tolist())}")
        return tendencia_sectores
    
    def analizar_areas_conocimiento(self):
        """Analiza tendencias por áreas de conocimiento"""
        print("\n🤖 Agente de Tendencias: Analizando áreas de conocimiento...")
        
        tendencia_areas = self.datos.groupby('AREA_CONOCIMIENTO')['MATRICULA'].sum().sort_values(ascending=False)
        
        self.tendencias['areas'] = tendencia_areas
        print(f"✓ Analizadas {len(tendencia_areas)} áreas de conocimiento")
        return tendencia_areas


class AgenteVisualizacion:
    """
    Agente de IA especializado en generación de visualizaciones.
    Crea gráficas profesionales para presentar los resultados del análisis.
    """
    
    def __init__(self, datos, tendencias):
        self.datos = datos
        self.tendencias = tendencias
        self.figuras = []
        
    def crear_grafica_tendencia_temporal(self):
        """Crea gráfica de tendencia temporal de matrícula"""
        print("\n🤖 Agente de Visualización: Creando gráfica de tendencia temporal...")
        
        fig, ax = plt.subplots(figsize=(12, 6))
        
        tendencia = self.tendencias.get('temporal')
        if tendencia is not None and len(tendencia) > 0:
            ax.plot(tendencia.index, tendencia.values, marker='o', linewidth=2, markersize=8)
            ax.set_xlabel('Período', fontsize=12, fontweight='bold')
            ax.set_ylabel('Matrícula Total', fontsize=12, fontweight='bold')
            ax.set_title('Tendencia de Matrícula a lo Largo del Tiempo', fontsize=14, fontweight='bold')
            ax.grid(True, alpha=0.3)
            plt.xticks(rotation=45, ha='right')
            plt.tight_layout()
            
            # Guardar figura
            plt.savefig('tendencia_temporal_matricula.png', dpi=300, bbox_inches='tight')
            self.figuras.append('tendencia_temporal_matricula.png')
            print("✓ Gráfica guardada: tendencia_temporal_matricula.png")
        
        plt.close()
    
    def crear_grafica_top_instituciones(self, top_n=10):
        """Crea gráfica de top instituciones por matrícula"""
        print("\n🤖 Agente de Visualización: Creando gráfica de top instituciones...")
        
        fig, ax = plt.subplots(figsize=(12, 8))
        
        tendencia = self.tendencias.get('institucional')
        if tendencia is not None and len(tendencia) > 0:
            top_instituciones = tendencia.head(top_n)
            
            colors = sns.color_palette("viridis", len(top_instituciones))
            ax.barh(range(len(top_instituciones)), top_instituciones.values, color=colors)
            ax.set_yticks(range(len(top_instituciones)))
            ax.set_yticklabels([inst[:50] + '...' if len(inst) > 50 else inst 
                                for inst in top_instituciones.index], fontsize=10)
            ax.set_xlabel('Matrícula Total', fontsize=12, fontweight='bold')
            ax.set_title(f'Top {top_n} Instituciones por Matrícula Total', fontsize=14, fontweight='bold')
            ax.invert_yaxis()
            plt.tight_layout()
            
            # Guardar figura
            plt.savefig('top_instituciones_matricula.png', dpi=300, bbox_inches='tight')
            self.figuras.append('top_instituciones_matricula.png')
            print("✓ Gráfica guardada: top_instituciones_matricula.png")
        
        plt.close()
    
    def crear_grafica_programas(self):
        """Crea gráfica de programas académicos"""
        print("\n🤖 Agente de Visualización: Creando gráfica de programas académicos...")
        
        fig, ax = plt.subplots(figsize=(10, 6))
        
        tendencia = self.tendencias.get('programas')
        if tendencia is not None and len(tendencia) > 0:
            colors = sns.color_palette("Set2", len(tendencia))
            ax.bar(range(len(tendencia)), tendencia.values, color=colors)
            ax.set_xticks(range(len(tendencia)))
            ax.set_xticklabels([prog[:30] + '...' if len(prog) > 30 else prog 
                                for prog in tendencia.index], rotation=45, ha='right', fontsize=9)
            ax.set_ylabel('Matrícula Total', fontsize=12, fontweight='bold')
            ax.set_title('Distribución de Matrícula por Programa Académico', fontsize=14, fontweight='bold')
            ax.grid(True, alpha=0.3, axis='y')
            plt.tight_layout()
            
            # Guardar figura
            plt.savefig('distribucion_programas.png', dpi=300, bbox_inches='tight')
            self.figuras.append('distribucion_programas.png')
            print("✓ Gráfica guardada: distribucion_programas.png")
        
        plt.close()
    
    def crear_grafica_sectores(self):
        """Crea gráfica de distribución por sectores"""
        print("\n🤖 Agente de Visualización: Creando gráfica de sectores...")
        
        fig, ax = plt.subplots(figsize=(10, 8))
        
        tendencia = self.tendencias.get('sectores')
        if tendencia is not None and len(tendencia) > 0:
            colors = ['#ff9999', '#66b3ff']
            explode = [0.05] * len(tendencia)
            
            ax.pie(tendencia.values, labels=tendencia.index, autopct='%1.1f%%',
                   startangle=90, colors=colors, explode=explode,
                   textprops={'fontsize': 12, 'fontweight': 'bold'})
            ax.set_title('Distribución de Matrícula por Sector (Público vs Privado)', 
                        fontsize=14, fontweight='bold', pad=20)
            plt.tight_layout()
            
            # Guardar figura
            plt.savefig('distribucion_sectores.png', dpi=300, bbox_inches='tight')
            self.figuras.append('distribucion_sectores.png')
            print("✓ Gráfica guardada: distribucion_sectores.png")
        
        plt.close()
    
    def crear_grafica_areas_conocimiento(self):
        """Crea gráfica de áreas de conocimiento"""
        print("\n🤖 Agente de Visualización: Creando gráfica de áreas de conocimiento...")
        
        fig, ax = plt.subplots(figsize=(10, 6))
        
        tendencia = self.tendencias.get('areas')
        if tendencia is not None and len(tendencia) > 0:
            colors = sns.color_palette("coolwarm", len(tendencia))
            bars = ax.bar(range(len(tendencia)), tendencia.values, color=colors)
            ax.set_xticks(range(len(tendencia)))
            ax.set_xticklabels([area[:40] + '...' if len(area) > 40 else area 
                                for area in tendencia.index], rotation=45, ha='right', fontsize=10)
            ax.set_ylabel('Matrícula Total', fontsize=12, fontweight='bold')
            ax.set_title('Distribución de Matrícula por Área de Conocimiento', 
                        fontsize=14, fontweight='bold')
            ax.grid(True, alpha=0.3, axis='y')
            
            # Añadir valores en las barras
            for bar in bars:
                height = bar.get_height()
                ax.text(bar.get_x() + bar.get_width()/2., height,
                       f'{int(height):,}',
                       ha='center', va='bottom', fontsize=9)
            
            plt.tight_layout()
            
            # Guardar figura
            plt.savefig('distribucion_areas_conocimiento.png', dpi=300, bbox_inches='tight')
            self.figuras.append('distribucion_areas_conocimiento.png')
            print("✓ Gráfica guardada: distribucion_areas_conocimiento.png")
        
        plt.close()
    
    def crear_dashboard_completo(self):
        """Crea un dashboard con múltiples gráficas"""
        print("\n🤖 Agente de Visualización: Creando dashboard completo...")
        
        fig = plt.figure(figsize=(16, 12))
        gs = fig.add_gridspec(3, 2, hspace=0.3, wspace=0.3)
        
        # Gráfica 1: Tendencia temporal
        ax1 = fig.add_subplot(gs[0, :])
        tendencia = self.tendencias.get('temporal')
        if tendencia is not None and len(tendencia) > 0:
            ax1.plot(tendencia.index, tendencia.values, marker='o', linewidth=2, markersize=6, color='#2E86AB')
            ax1.set_xlabel('Período', fontsize=10, fontweight='bold')
            ax1.set_ylabel('Matrícula Total', fontsize=10, fontweight='bold')
            ax1.set_title('Evolución Temporal de Matrícula', fontsize=12, fontweight='bold')
            ax1.grid(True, alpha=0.3)
            plt.setp(ax1.xaxis.get_majorticklabels(), rotation=45, ha='right')
        
        # Gráfica 2: Top 5 instituciones
        ax2 = fig.add_subplot(gs[1, 0])
        tendencia_inst = self.tendencias.get('institucional')
        if tendencia_inst is not None and len(tendencia_inst) > 0:
            top5_inst = tendencia_inst.head(5)
            colors = sns.color_palette("viridis", 5)
            ax2.barh(range(len(top5_inst)), top5_inst.values, color=colors)
            ax2.set_yticks(range(len(top5_inst)))
            ax2.set_yticklabels([inst[:30] for inst in top5_inst.index], fontsize=8)
            ax2.set_xlabel('Matrícula', fontsize=10, fontweight='bold')
            ax2.set_title('Top 5 Instituciones', fontsize=12, fontweight='bold')
            ax2.invert_yaxis()
        
        # Gráfica 3: Sectores
        ax3 = fig.add_subplot(gs[1, 1])
        tendencia_sect = self.tendencias.get('sectores')
        if tendencia_sect is not None and len(tendencia_sect) > 0:
            colors = ['#ff9999', '#66b3ff']
            ax3.pie(tendencia_sect.values, labels=tendencia_sect.index, autopct='%1.1f%%',
                   startangle=90, colors=colors, textprops={'fontsize': 9})
            ax3.set_title('Distribución por Sector', fontsize=12, fontweight='bold')
        
        # Gráfica 4: Programas
        ax4 = fig.add_subplot(gs[2, 0])
        tendencia_prog = self.tendencias.get('programas')
        if tendencia_prog is not None and len(tendencia_prog) > 0:
            colors = sns.color_palette("Set2", len(tendencia_prog))
            ax4.bar(range(len(tendencia_prog)), tendencia_prog.values, color=colors)
            ax4.set_xticks(range(len(tendencia_prog)))
            ax4.set_xticklabels([prog[:15] for prog in tendencia_prog.index], 
                               rotation=45, ha='right', fontsize=7)
            ax4.set_ylabel('Matrícula', fontsize=10, fontweight='bold')
            ax4.set_title('Programas Académicos', fontsize=12, fontweight='bold')
            ax4.grid(True, alpha=0.3, axis='y')
        
        # Gráfica 5: Áreas de conocimiento
        ax5 = fig.add_subplot(gs[2, 1])
        tendencia_areas = self.tendencias.get('areas')
        if tendencia_areas is not None and len(tendencia_areas) > 0:
            colors = sns.color_palette("coolwarm", len(tendencia_areas))
            ax5.bar(range(len(tendencia_areas)), tendencia_areas.values, color=colors)
            ax5.set_xticks(range(len(tendencia_areas)))
            ax5.set_xticklabels([area[:20] for area in tendencia_areas.index], 
                               rotation=45, ha='right', fontsize=7)
            ax5.set_ylabel('Matrícula', fontsize=10, fontweight='bold')
            ax5.set_title('Áreas de Conocimiento', fontsize=12, fontweight='bold')
            ax5.grid(True, alpha=0.3, axis='y')
        
        plt.suptitle('Dashboard de Análisis de Tendencias de Mercado Académico', 
                    fontsize=16, fontweight='bold', y=0.995)
        
        # Guardar dashboard
        plt.savefig('dashboard_completo.png', dpi=300, bbox_inches='tight')
        self.figuras.append('dashboard_completo.png')
        print("✓ Dashboard guardado: dashboard_completo.png")
        plt.close()


class AgenteRecomendaciones:
    """
    Agente de IA especializado en generar insights y recomendaciones.
    Analiza los resultados y proporciona conclusiones estratégicas.
    """
    
    def __init__(self, datos, tendencias):
        self.datos = datos
        self.tendencias = tendencias
        self.recomendaciones = []
        
    def generar_insights(self):
        """Genera insights basados en el análisis de datos"""
        print("\n🤖 Agente de Recomendaciones: Generando insights...")
        
        insights = []
        
        # Insight 1: Institución líder
        if 'institucional' in self.tendencias and len(self.tendencias['institucional']) > 0:
            institucion_lider = self.tendencias['institucional'].index[0]
            matricula_lider = self.tendencias['institucional'].values[0]
            insights.append(f"📌 La institución líder en matrícula es '{institucion_lider}' "
                          f"con {int(matricula_lider):,} estudiantes matriculados.")
        
        # Insight 2: Programa más popular
        if 'programas' in self.tendencias and len(self.tendencias['programas']) > 0:
            programa_popular = self.tendencias['programas'].index[0]
            matricula_programa = self.tendencias['programas'].values[0]
            insights.append(f"📌 El programa académico más popular es '{programa_popular}' "
                          f"con {int(matricula_programa):,} estudiantes.")
        
        # Insight 3: Distribución sectorial
        if 'sectores' in self.tendencias and len(self.tendencias['sectores']) > 0:
            total_matricula = self.tendencias['sectores'].sum()
            for sector, matricula in self.tendencias['sectores'].items():
                porcentaje = (matricula / total_matricula) * 100
                insights.append(f"📌 El sector {sector} representa el {porcentaje:.1f}% "
                              f"de la matrícula total.")
        
        # Insight 4: Área de conocimiento dominante
        if 'areas' in self.tendencias and len(self.tendencias['areas']) > 0:
            area_dominante = self.tendencias['areas'].index[0]
            matricula_area = self.tendencias['areas'].values[0]
            insights.append(f"📌 El área de conocimiento dominante es '{area_dominante}' "
                          f"con {int(matricula_area):,} estudiantes.")
        
        return insights
    
    def generar_recomendaciones_estrategicas(self):
        """Genera recomendaciones estratégicas"""
        print("\n🤖 Agente de Recomendaciones: Generando recomendaciones estratégicas...")
        
        recomendaciones = []
        
        # Recomendación 1: Oportunidades de mercado
        if 'programas' in self.tendencias and len(self.tendencias['programas']) > 0:
            recomendaciones.append("💡 Oportunidad de Mercado: Los programas con mayor demanda "
                                 "representan áreas de alta oportunidad para inversión y expansión.")
        
        # Recomendación 2: Análisis competitivo
        if 'institucional' in self.tendencias and len(self.tendencias['institucional']) > 0:
            recomendaciones.append("💡 Análisis Competitivo: Las instituciones líderes marcan "
                                 "las tendencias del mercado. Estudiar sus estrategias puede "
                                 "proporcionar ventajas competitivas.")
        
        # Recomendación 3: Diversificación
        if 'areas' in self.tendencias and len(self.tendencias['areas']) > 1:
            recomendaciones.append("💡 Diversificación: Considerar la diversificación hacia "
                                 "múltiples áreas de conocimiento para mitigar riesgos y "
                                 "capturar diferentes segmentos de mercado.")
        
        # Recomendación 4: Sector público vs privado
        if 'sectores' in self.tendencias and len(self.tendencias['sectores']) > 0:
            recomendaciones.append("💡 Estrategia Sectorial: Analizar las diferencias entre "
                                 "sectores público y privado para identificar nichos "
                                 "específicos de oportunidad.")
        
        return recomendaciones
    
    def generar_reporte_completo(self, insights, recomendaciones):
        """Genera un reporte completo en formato texto"""
        print("\n🤖 Agente de Recomendaciones: Generando reporte completo...")
        
        reporte = []
        reporte.append("=" * 80)
        reporte.append("REPORTE DE ANÁLISIS DE TENDENCIAS DE MERCADO ACADÉMICO")
        reporte.append("=" * 80)
        reporte.append(f"Fecha de generación: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
        reporte.append("")
        
        reporte.append("\n" + "=" * 80)
        reporte.append("INSIGHTS PRINCIPALES")
        reporte.append("=" * 80)
        for insight in insights:
            reporte.append(insight)
        
        reporte.append("\n" + "=" * 80)
        reporte.append("RECOMENDACIONES ESTRATÉGICAS")
        reporte.append("=" * 80)
        for rec in recomendaciones:
            reporte.append(rec)
        
        reporte.append("\n" + "=" * 80)
        reporte.append("CONCLUSIÓN")
        reporte.append("=" * 80)
        reporte.append("Este análisis proporciona una visión integral de las tendencias de mercado")
        reporte.append("en el sector académico. Las visualizaciones generadas permiten identificar")
        reporte.append("patrones clave, oportunidades de crecimiento y áreas de mejora estratégica.")
        reporte.append("=" * 80)
        
        # Guardar reporte
        with open('reporte_analisis_tendencias.txt', 'w', encoding='utf-8') as f:
            f.write('\n'.join(reporte))
        
        print("✓ Reporte guardado: reporte_analisis_tendencias.txt")
        
        # También imprimir en consola
        print('\n'.join(reporte))

from pptx import Presentation
from pptx.util import Inches

class AgentePresentacion:
    def __init__(self, resumen, tendencias, insights, recomendaciones, imagenes):
        self.resumen = resumen
        self.tendencias = tendencias
        self.insights = insights
        self.recomendaciones = recomendaciones
        self.imagenes = imagenes  # lista de rutas de imágenes (png)

    def generar_presentacion(self, nombre_archivo="presentacion_analisis_tendencias.pptx"):
        prs = Presentation()

        # 1. Diapositiva de título
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Análisis de Tendencias de Mercado Académico"
        slide.placeholders[1].text = "Reporte ejecutivo generado por agentes de IA"

        # 2. Diapositiva de resumen
        summary_slide = prs.slides.add_slide(prs.slide_layouts[1])
        summary_slide.shapes.title.text = "Resumen de datos"
        cuerpo = summary_slide.placeholders[1].text_frame
        cuerpo.clear()
        cuerpo.add_paragraph().text = f"Total de registros: {self.resumen['total_registros']:,}"
        cuerpo.add_paragraph().text = f"Instituciones únicas: {self.resumen['instituciones']:,}"
        cuerpo.add_paragraph().text = f"Programas académicos únicos: {self.resumen['programas']:,}"
        cuerpo.add_paragraph().text = f"Períodos analizados: {self.resumen['periodos']:,}"

        # 3. Diapositivas con gráficos
        for ruta in self.imagenes:
            titulo = ruta.replace(".png", "").replace("_", " ").title()
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = titulo
            slide.shapes.add_picture(ruta, Inches(0.5), Inches(1.5), width=Inches(9))

        # 4. Diapositiva de insights y recomendaciones
        info_slide = prs.slides.add_slide(prs.slide_layouts[1])
        info_slide.shapes.title.text = "Insights y recomendaciones"
        txt_frame = info_slide.placeholders[1].text_frame
        txt_frame.clear()
        txt_frame.add_paragraph().text = "📌 Principales insights:"
        for insight in self.insights:
            p = txt_frame.add_paragraph()
            p.text = insight
            p.level = 1
        txt_frame.add_paragraph().text = "💡 Recomendaciones estratégicas:"
        for rec in self.recomendaciones:
            p = txt_frame.add_paragraph()
            p.text = rec
            p.level = 1

        # 5. Guardar presentación
        prs.save(nombre_archivo)
        print(f"✓ Presentación guardada: {nombre_archivo}")


def main():
    """
    Función principal que coordina todos los agentes de IA para realizar
    el análisis completo de tendencias de mercado.
    """
    print("\n" + "=" * 80)
    print("SISTEMA DE ANÁLISIS DE TENDENCIAS DE MERCADO CON AGENTES DE IA")
    print("=" * 80)
    print("Iniciando análisis de datos académicos...")
    print("=" * 80 + "\n")
    
    archivo_excel = 'programas_equivalentes.xlsx'
    
    # Fase 1: Análisis y preprocesamiento de datos
    print("\n📋 FASE 1: ANÁLISIS Y PREPROCESAMIENTO DE DATOS")
    print("-" * 80)
    agente_datos = AgenteAnalisisDatos(archivo_excel)
    datos = agente_datos.cargar_datos()
    datos = agente_datos.preprocesar_datos()
    resumen = agente_datos.obtener_resumen_estadistico()
    
    # Fase 2: Identificación de tendencias
    print("\n📊 FASE 2: IDENTIFICACIÓN DE TENDENCIAS DE MERCADO")
    print("-" * 80)
    agente_tendencias = AgenteTendenciasMercado(datos)
    agente_tendencias.analizar_tendencia_matricula_temporal()
    agente_tendencias.analizar_tendencia_por_institucion()
    agente_tendencias.analizar_tendencia_por_programa()
    agente_tendencias.analizar_distribucion_sectores()
    agente_tendencias.analizar_areas_conocimiento()
    
    # Fase 3: Generación de visualizaciones
    print("\n📈 FASE 3: GENERACIÓN DE VISUALIZACIONES")
    print("-" * 80)
    agente_viz = AgenteVisualizacion(datos, agente_tendencias.tendencias)
    agente_viz.crear_grafica_tendencia_temporal()
    agente_viz.crear_grafica_top_instituciones(top_n=10)
    agente_viz.crear_grafica_programas()
    agente_viz.crear_grafica_sectores()
    agente_viz.crear_grafica_areas_conocimiento()
    agente_viz.crear_dashboard_completo()
    
    # Fase 4: Generación de insights y recomendaciones
    print("\n💡 FASE 4: GENERACIÓN DE INSIGHTS Y RECOMENDACIONES")
    print("-" * 80)
    agente_recom = AgenteRecomendaciones(datos, agente_tendencias.tendencias)
    insights = agente_recom.generar_insights()
    recomendaciones = agente_recom.generar_recomendaciones_estrategicas()
    agente_recom.generar_reporte_completo(insights, recomendaciones)
    
    # Resumen final
    print("\n" + "=" * 80)
    print("✅ ANÁLISIS COMPLETADO EXITOSAMENTE")
    print("=" * 80)
    print(f"\nArchivos generados:")
    for i, figura in enumerate(agente_viz.figuras, 1):
        print(f"  {i}. {figura}")
    print(f"  {len(agente_viz.figuras) + 1}. reporte_analisis_tendencias.txt")
    print("\n" + "=" * 80 + "\n")

    # ... dentro de la función main()
    agente_presentacion = AgentePresentacion(
    resumen=resumen,
    tendencias=agente_tendencias.tendencias,
    insights=insights,
    recomendaciones=recomendaciones,
    imagenes=agente_viz.figuras)
    agente_presentacion.generar_presentacion("presentacion_analisis_tendencias.pptx")


if __name__ == "__main__":
    main()