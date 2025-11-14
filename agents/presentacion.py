from pptx import Presentation
from pptx.util import Inches

class AgentePresentacion:
    def __init__(self, resumen, tendencias, insights, recomendaciones, imagenes):
        self.resumen = resumen
        self.tendencias = tendencias
        self.insights = insights
        self.recomendaciones = recomendaciones
        self.imagenes = imagenes  # lista de rutas de imágenes (png)

    def generar_presentacion(self, nombre_archivo="presentacion_analisis_tendencias.pptx"):
        prs = Presentation()

        # 1. Diapositiva de título
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Análisis de Tendencias de Mercado Académico"
        slide.placeholders[1].text = "Reporte ejecutivo generado por agentes de IA"

        # 2. Diapositiva de resumen
        summary_slide = prs.slides.add_slide(prs.slide_layouts[1])
        summary_slide.shapes.title.text = "Resumen de datos"
        cuerpo = summary_slide.placeholders[1].text_frame
        cuerpo.clear()
        cuerpo.add_paragraph().text = f"Total de registros: {self.resumen['total_registros']:,}"
        cuerpo.add_paragraph().text = f"Instituciones únicas: {self.resumen['instituciones']:,}"
        cuerpo.add_paragraph().text = f"Programas académicos únicos: {self.resumen['programas']:,}"
        cuerpo.add_paragraph().text = f"Períodos analizados: {self.resumen['periodos']:,}"

        # 3. Diapositivas con gráficos
        for ruta in self.imagenes:
            titulo = ruta.replace(".png", "").replace("_", " ").title()
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = titulo
            slide.shapes.add_picture(ruta, Inches(0.5), Inches(1.5), width=Inches(9))

        # 4. Diapositiva de insights y recomendaciones
        info_slide = prs.slides.add_slide(prs.slide_layouts[1])
        info_slide.shapes.title.text = "Insights y recomendaciones"
        txt_frame = info_slide.placeholders[1].text_frame
        txt_frame.clear()
        txt_frame.add_paragraph().text = "📌 Principales insights:"
        for insight in self.insights:
            p = txt_frame.add_paragraph()
            p.text = insight
            p.level = 1
        txt_frame.add_paragraph().text = "💡 Recomendaciones estratégicas:"
        for rec in self.recomendaciones:
            p = txt_frame.add_paragraph()
            p.text = rec
            p.level = 1

        # 5. Guardar presentación
        prs.save(nombre_archivo)
        print(f"Presentación guardada: {nombre_archivo}")
